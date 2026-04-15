"""
Hybrid retrieval system with:
- Glossary prioritisation (exact/fuzzy)
- Email boosting for sender queries
- Glossary downweighting for broad queries
- RRF fusion + caching
"""

import os
import json
import re
from pathlib import Path
from typing import List, Dict, Tuple, Any, Optional

# Optional imports
try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from docx import Document
except ImportError:
    Document = None
try:
    from email import policy
    from email.parser import BytesParser
except ImportError:
    BytesParser = None
try:
    from rank_bm25 import BM25Okapi
except ImportError:
    BM25Okapi = None
try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    SentenceTransformer = None
import numpy as np

_CORPUS: List[Dict] = []
_BM25: Any = None
_EMBEDDINGS: Optional[np.ndarray] = None
_EMBEDDER: Any = None
_GLOSSARY: Dict[str, str] = {}
MIN_CHUNK_CHARS = 50


# ----------------------------------------------------------------------
# Text extraction (returns (text, glossary_dict))
# ----------------------------------------------------------------------
def _extract_from_pdf(path: Path) -> Tuple[str, Dict]:
    if pdfplumber is None:
        return "", {}
    text = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text.append(t)
    except Exception:
        pass
    return "\n".join(text), {}

def _extract_from_pptx(path: Path) -> Tuple[str, Dict]:
    if Presentation is None:
        return "", {}
    text = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
    except Exception:
        pass
    return "\n".join(text), {}

def _extract_from_docx(path: Path) -> Tuple[str, Dict]:
    if Document is None:
        return "", {}
    text = []
    try:
        doc = Document(path)
        for para in doc.paragraphs:
            if para.text:
                text.append(para.text)
    except Exception:
        pass
    return "\n".join(text), {}

def _extract_from_eml(path: Path) -> Tuple[str, Dict]:
    if BytesParser is None:
        return "", {}
    text = []
    try:
        with open(path, "rb") as f:
            msg = BytesParser(policy=policy.default).parse(f)
        subject = msg.get("Subject", "")
        if subject:
            text.append(f"Subject: {subject}")
        body = None
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_content()
                    break
        else:
            body = msg.get_content()
        if body:
            text.append(body)
    except Exception:
        pass
    return "\n".join(text), {}

def _normalize_glossary_key(value: str) -> str:
    normalized = value.strip().lower().replace("-", " ").replace("/", " ")
    normalized = normalized.replace("_", " ")
    normalized = re.sub(r"\s+", " ", normalized)
    return normalized

def _is_placeholder_definition(key: str, definition: str) -> bool:
    key_normalized = _normalize_glossary_key(key)
    definition_normalized = _normalize_glossary_key(definition)
    return not definition_normalized or definition_normalized == key_normalized

def _store_glossary_entry(glossary: Dict[str, str], full_lines: List[str], key: str, definition: str) -> None:
    raw_key = key.strip().lower()
    cleaned_definition = definition.strip()
    if not raw_key or not cleaned_definition:
        return
    glossary[raw_key] = cleaned_definition
    full_lines.append(f"{key}: {cleaned_definition}")

def _extract_glossary_entries(data: Any, glossary: Dict[str, str], full_lines: List[str]) -> None:
    if not isinstance(data, dict):
        return

    definition = data.get("definition")
    if isinstance(definition, str):
        for key, value in data.items():
            if key == "definition" or not isinstance(value, dict):
                continue
            nested_definition = value.get("definition")
            if isinstance(nested_definition, str):
                _store_glossary_entry(glossary, full_lines, key, nested_definition)
        return

    for key, value in data.items():
        if isinstance(value, dict):
            child_definition = value.get("definition")
            if isinstance(child_definition, str):
                _store_glossary_entry(glossary, full_lines, key, child_definition)
            _extract_glossary_entries(value, glossary, full_lines)

def _extract_from_json_glossary(path: Path) -> Tuple[str, Dict[str, str]]:
    glossary: Dict[str, str] = {}
    full_lines: List[str] = []
    encodings = ["utf-8", "latin-1"]
    data = None
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc) as f:
                data = json.load(f)
            break
        except (UnicodeDecodeError, json.JSONDecodeError):
            continue
    if data is None:
        print(f"Warning: could not decode {path} with utf-8 or latin-1")
        return "", glossary

    _extract_glossary_entries(data, glossary, full_lines)
    return "\n".join(full_lines), glossary

def _extract_text_from_file(file_path: Path) -> Tuple[str, Dict[str, str]]:
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return _extract_from_pdf(file_path)
    elif ext == ".pptx":
        return _extract_from_pptx(file_path)
    elif ext == ".docx":
        return _extract_from_docx(file_path)
    elif ext == ".eml":
        return _extract_from_eml(file_path)
    elif ext == ".json" and "glossary" in file_path.name.lower():
        return _extract_from_json_glossary(file_path)
    return "", {}


# ----------------------------------------------------------------------
# Chunking
# ----------------------------------------------------------------------
def _chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    if not text:
        return []
    words = text.split()
    if len(words) <= chunk_size:
        return [text] if len(text.strip()) >= MIN_CHUNK_CHARS else []
    chunks = []
    step = chunk_size - overlap
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i+chunk_size])
        if len(chunk.strip()) >= MIN_CHUNK_CHARS:
            chunks.append(chunk)
        if i + chunk_size >= len(words):
            break
    return chunks


# ----------------------------------------------------------------------
# Build corpus
# ----------------------------------------------------------------------
def _build_corpus(data_dir: str = "test_files") -> None:
    global _CORPUS, _GLOSSARY, _BM25, _EMBEDDINGS, _EMBEDDER
    if _CORPUS:
        return

    data_path = Path(data_dir)
    if not data_path.exists():
        module_data_path = Path(__file__).resolve().parent / "test_files"
        if module_data_path.exists():
            data_path = module_data_path
        elif (Path(".") / "test_files").exists():
            data_path = Path(".") / "test_files"
        else:
            data_path = Path(__file__).resolve().parent

    all_chunks = []
    glossary_global = {}
    chunk_id = 0

    for file_path in data_path.glob("*"):
        if file_path.is_dir():
            continue
        text, file_glossary = _extract_text_from_file(file_path)
        if file_glossary:
            glossary_global.update(file_glossary)
        if not text:
            continue
        chunks = _chunk_text(text)
        for chunk in chunks:
            all_chunks.append({
                "id": chunk_id,
                "text": chunk,
                "source": file_path.name,
            })
            chunk_id += 1

    _CORPUS = all_chunks
    _GLOSSARY = glossary_global

    print(f"[INFO] Loaded glossary with {len(_GLOSSARY)} entries. Sample: {list(_GLOSSARY.keys())[:5]}")

    if BM25Okapi is not None and _CORPUS:
        tokenized = [doc["text"].lower().split() for doc in _CORPUS]
        _BM25 = BM25Okapi(tokenized)
    if SentenceTransformer is not None and _CORPUS:
        _EMBEDDER = SentenceTransformer("all-MiniLM-L6-v2")
        texts = [doc["text"] for doc in _CORPUS]
        _EMBEDDINGS = _EMBEDDER.encode(texts, show_progress_bar=False)


# ----------------------------------------------------------------------
# Retrieval functions
# ----------------------------------------------------------------------
def _dense_retrieve(query: str, top_k: int = 20) -> List[Tuple[int, float]]:
    if _EMBEDDINGS is None or _EMBEDDER is None or not _CORPUS:
        return []
    qvec = _EMBEDDER.encode([query])[0]
    norms = np.linalg.norm(_EMBEDDINGS, axis=1)
    qnorm = np.linalg.norm(qvec)
    if qnorm == 0:
        return []
    sim = np.dot(_EMBEDDINGS, qvec) / (norms * qnorm + 1e-8)
    top = np.argsort(sim)[-top_k:][::-1]
    return [(int(idx), float(sim[idx])) for idx in top]

def _bm25_retrieve(query: str, top_k: int = 20) -> List[Tuple[int, float]]:
    if _BM25 is None or not _CORPUS:
        return []
    tokens = query.lower().split()
    scores = _BM25.get_scores(tokens)
    top = np.argsort(scores)[-top_k:][::-1]
    return [(int(idx), float(scores[idx])) for idx in top]

def _apply_boosts(query: str, 
                  bm25_results: List[Tuple[int, float]], 
                  dense_results: List[Tuple[int, float]]) -> Tuple[List[Tuple[int, float]], List[Tuple[int, float]]]:
    """
    Apply source-specific boosts based on query content.
    - Email boost: if query contains email-related terms, multiply .eml scores by 1.5
    - Glossary downweight: for non-variable queries, multiply glossary chunks by 0.5
    """
    query_lower = query.lower()
    email_keywords = ["email", "sent", "sender", "who sent", "from:", "subject:"]
    is_email_query = any(kw in query_lower for kw in email_keywords)
    
    has_variable_pattern = bool(re.search(r"\b[a-z]+_[a-z_]+\b", query_lower))
    is_broad_query = not has_variable_pattern
    
    def boost_list(results):
        new_results = []
        for doc_id, score in results:
            source = _CORPUS[doc_id]["source"]
            if is_email_query and source.endswith(".eml"):
                score *= 1.5
            if is_broad_query and source == "glossary_partial.json":
                score *= 0.5
            new_results.append((doc_id, score))
        new_results.sort(key=lambda x: x[1], reverse=True)
        return new_results
    
    return boost_list(bm25_results), boost_list(dense_results)

def _is_narrow_query(query: str) -> bool:
    query_lower = query.lower()
    if re.search(r"\b(?:frm|var|q)_[a-z0-9_]{3,}\b", query_lower):
        return True
    if re.search(r"\b[a-z]+_[a-z0-9_]{4,}\b", query_lower):
        return True
    return bool(re.search(r"(?:what does|define|meaning of|measure|scale|represent)\s+`?[a-z_]+`?", query_lower))

def _rrf_fusion(lists: List[List[Tuple[int, float]]], k: int = 60) -> List[Tuple[int, float]]:
    rrf = {}
    for lst in lists:
        for rank, (doc_id, _) in enumerate(lst):
            rrf[doc_id] = rrf.get(doc_id, 0) + 1.0 / (k + rank + 1)
    return sorted(rrf.items(), key=lambda x: x[1], reverse=True)


# ----------------------------------------------------------------------
# Glossary lookup (improved with prefix stripping)
# ----------------------------------------------------------------------
def _try_glossary_lookup(query: str) -> Optional[List[Dict]]:
    if not _GLOSSARY or not _is_narrow_query(query):
        return None

    query_lower = query.lower()
    candidates: List[str] = []

    matches = re.findall(r"`([a-z0-9_]+)`", query_lower)
    candidates.extend(matches)
    words = re.findall(r"\b[a-z][a-z0-9_]{3,}\b", query_lower)
    for word in words:
        if "_" in word and len(word) <= 50:
            candidates.append(word)
    measure_match = re.search(r"(?:what does|define|meaning of)\s+`?([a-z_]+)`?", query_lower)
    if measure_match:
        candidates.append(measure_match.group(1))

    normalized_candidates: List[str] = []
    seen = set()
    for candidate in candidates:
        variants = [
            candidate.strip().lower(),
            _normalize_glossary_key(candidate),
        ]
        for prefix in ["frm_", "var_", "q_"]:
            if candidate.startswith(prefix):
                stripped = candidate[len(prefix):]
                variants.append(stripped)
                variants.append(_normalize_glossary_key(stripped))
        for variant in variants:
            if variant and variant not in seen:
                normalized_candidates.append(variant)
                seen.add(variant)

    if not normalized_candidates and not _is_narrow_query(query):
        return None

    if not candidates:
        fallback_tokens = re.findall(r"\b[a-z][a-z0-9_]{4,}\b", query_lower)
        for token in fallback_tokens:
            normalized = _normalize_glossary_key(token)
            if normalized and normalized not in seen:
                normalized_candidates.append(normalized)
                seen.add(normalized)

    best_match = None
    best_score = -1
    for cand in normalized_candidates:
        cand_normalized = _normalize_glossary_key(cand)
        stripped_variants = []
        for prefix in ["frm_", "var_", "q_"]:
            if cand.startswith(prefix):
                stripped = cand[len(prefix):]
                stripped_variants.extend([stripped, _normalize_glossary_key(stripped)])

        if cand in _GLOSSARY:
            definition = _GLOSSARY[cand]
            if not _is_placeholder_definition(cand, definition):
                return [{
                    "text": f"Variable '{cand}' measures: {definition}",
                    "source": "glossary",
                    "score": 1.0
                }]

        for key, desc in _GLOSSARY.items():
            key_normalized = _normalize_glossary_key(key)
            score = -1
            if cand_normalized == key_normalized:
                score = 300
            elif cand == key:
                score = 290
            elif any(variant and (key == variant or key_normalized == variant) for variant in stripped_variants):
                score = 260
            elif any(variant and (key.endswith(f"_{variant}") or key_normalized.endswith(f" {variant}")) for variant in stripped_variants):
                score = 250
            elif any(variant and variant in key_normalized for variant in stripped_variants):
                score = 225
            elif cand_normalized in key_normalized or key_normalized in cand_normalized:
                score = 200

            if score < 0:
                continue
            is_placeholder = _is_placeholder_definition(key, desc)
            if is_placeholder:
                score -= 120
            else:
                score += 15
            score += min(len(desc.strip()) // 40, 10)
            if score > best_score:
                best_score = score
                best_match = (key, desc, cand, score)

    if best_match is not None:
        key, desc, cand, score = best_match
        exactish = score >= 300 or _normalize_glossary_key(key) == _normalize_glossary_key(cand)
        confidence = 1.0 if exactish else 0.95
        if exactish:
            text = f"Variable '{key}' measures: {desc}"
        else:
            text = f"Variable '{key}' (close match to '{cand}'): {desc}"
        return [{
            "text": text,
            "source": "glossary",
            "score": confidence
        }]
    return None


# ----------------------------------------------------------------------
# Main retrieve
# ----------------------------------------------------------------------
def retrieve(query: str) -> List[Dict]:
    _build_corpus()

    glossary_result = _try_glossary_lookup(query)
    if glossary_result:
        # Get hybrid results and prepend glossary hit
        bm25_res = _bm25_retrieve(query, top_k=20)
        dense_res = _dense_retrieve(query, top_k=20)

        if bm25_res or dense_res:
            bm25_boosted, dense_boosted = _apply_boosts(query, bm25_res, dense_res)
            fused = _rrf_fusion([bm25_boosted, dense_boosted])
            if fused:
                max_score = fused[0][1]
                min_score = fused[-1][1] if len(fused) > 1 else max_score
                if max_score == min_score:
                    norm = [1.0] * len(fused)
                else:
                    norm = [(s - min_score) / (max_score - min_score) for _, s in fused]

                remaining = []
                for i, (doc_id, _) in enumerate(fused[:4]):
                    chunk = _CORPUS[doc_id]
                    remaining.append({
                        "text": chunk["text"],
                        "source": chunk["source"],
                        "score": norm[i]
                    })
                return glossary_result + remaining
        # Fallback if no hybrid results
        return glossary_result

    bm25_res = _bm25_retrieve(query, top_k=20)
    dense_res = _dense_retrieve(query, top_k=20)

    if not bm25_res and not dense_res:
        terms = set(query.lower().split())
        candidates = []
        for idx, chunk in enumerate(_CORPUS):
            text_lower = chunk["text"].lower()
            score = sum(1 for t in terms if t in text_lower) / (len(terms) + 1e-6)
            if score > 0:
                candidates.append((idx, score))
        candidates.sort(key=lambda x: x[1], reverse=True)
        results = []
        for idx, score in candidates[:5]:
            results.append({
                "text": _CORPUS[idx]["text"],
                "source": _CORPUS[idx]["source"],
                "score": score
            })
        return results

    bm25_boosted, dense_boosted = _apply_boosts(query, bm25_res, dense_res)
    fused = _rrf_fusion([bm25_boosted, dense_boosted])
    if not fused:
        return []

    max_score = fused[0][1]
    min_score = fused[-1][1] if len(fused) > 1 else max_score
    if max_score == min_score:
        norm = [1.0] * len(fused)
    else:
        norm = [(s - min_score) / (max_score - min_score) for _, s in fused]

    results = []
    for i, (doc_id, _) in enumerate(fused[:5]):
        chunk = _CORPUS[doc_id]
        results.append({
            "text": chunk["text"],
            "source": chunk["source"],
            "score": norm[i]
        })
    return results
